"""Generate the 5-minute presentation deck (PPTX) for the churn MLOps project."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- palette ---
NAVY = RGBColor(0x14, 0x2A, 0x4A)
TEAL = RGBColor(0x14, 0x8F, 0x96)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x2B, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         bullet=False, space=8):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def title_bar(slide, text, kicker=None):
    """Top band with a slide title."""
    band = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame; tf.margin_left = Inches(0.5); tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    # teal accent stripe
    stripe = slide.shapes.add_shape(1, 0, Inches(1.15), SW, Inches(0.07))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = TEAL; stripe.line.fill.background()


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
tf = box(s, 1, 2.2, 11.3, 2.6)
para(tf, "Machine Learning Lifecycle Management with MLflow", 40, WHITE, bold=True, first=True)
para(tf, "Predicting E-Commerce Customer Churn — an end-to-end MLOps system", 22, TEAL, space=4)
tf2 = box(s, 1, 5.3, 11.3, 1.6)
para(tf2, "Asra Sarı  ·  Student No. 2101640", 20, WHITE, bold=True, first=True)
para(tf2, "AIN-3009  Delivering AI Applications with MLOps  ·  Bahçeşehir University", 16, LIGHT)
para(tf2, "github.com/asrasari/ecommerce-churn-mlops", 14, TEAL)

# ---------- Slide 2: Problem & Dataset ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE); title_bar(s, "The Problem & The Data")
tf = box(s, 0.7, 1.6, 7.3, 5.5)
para(tf, "Why churn?", 22, TEAL, bold=True, first=True)
para(tf, "Retaining a customer is far cheaper than acquiring a new one — predicting who will leave lets a business act early.", 18, DARK, space=14)
para(tf, "But the real focus: managing the ML lifecycle with MLflow", 22, TEAL, bold=True)
para(tf, "Track experiments · tune · register · deploy · monitor", 18, DARK, space=14)
para(tf, "Dataset", 22, TEAL, bold=True)
para(tf, "E-Commerce Customer Churn — 3,941 customers, 10 features", 18, DARK, bullet=True)
para(tf, "Target: Churn (yes/no), ~17% churn — a class imbalance", 18, DARK, bullet=True)
para(tf, "8 numeric + 2 categorical features; 576 missing values handled", 18, DARK, bullet=True)
# side card
card = s.shapes.add_shape(1, Inches(8.4), Inches(1.7), Inches(4.3), Inches(5.0))
card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = TEAL
ctf = card.text_frame; ctf.word_wrap = True; ctf.margin_left = Inches(0.3); ctf.margin_top = Inches(0.3)
para(ctf, "Example features", 18, NAVY, bold=True, first=True, space=10)
for feat in ["Tenure", "Satisfaction score", "Complaint filed?", "Days since last order",
             "Cashback amount", "Preferred order category", "Marital status"]:
    para(ctf, feat, 16, DARK, bullet=True, space=6)

# ---------- Slide 3: Architecture / lifecycle ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE); title_bar(s, "System Architecture — the MLflow Lifecycle")
stages = [
    ("1  Track", "Log every model run:\nparams, metrics,\nplots, model"),
    ("2  Tune", "Optuna search,\neach trial a\nnested run"),
    ("3  Register", "Best model →\nStaging →\nProduction"),
    ("4  Serve", "Load Production\nmodel, predict\nvia REST"),
    ("5  Monitor", "Live metrics +\nPSI drift over\nincoming batches"),
]
x = 0.55
for i, (h, d) in enumerate(stages):
    c = s.shapes.add_shape(1, Inches(x), Inches(2.2), Inches(2.25), Inches(2.7))
    c.fill.solid(); c.fill.fore_color.rgb = NAVY if i % 2 == 0 else TEAL
    c.line.fill.background()
    t = c.text_frame; t.word_wrap = True; t.margin_top = Inches(0.2)
    para(t, h, 20, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER, space=8)
    para(t, d, 13, WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        a = box(s, x + 2.27, 3.3, 0.45, 0.6)
        para(a, "→", 28, GREY, bold=True, first=True, align=PP_ALIGN.CENTER)
    x += 2.55
tf = box(s, 0.7, 5.4, 12, 1.5)
para(tf, "A local MLflow tracking server (port 8080) with a SQLite metadata database and a local "
         "artifact store ties it together. An Airflow DAG orchestrates the pipeline on a schedule.",
     17, DARK, first=True)
para(tf, "Code: focused Python modules — data_prep · train · tune · register · serve · monitor",
     15, GREY)

# ---------- Slide 4: Experiment tracking + results ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE); title_bar(s, "Experiment Tracking — Results")
tf = box(s, 0.7, 1.5, 6.0, 1.2)
para(tf, "Three models trained, each logged as an MLflow run and compared in the UI.",
     18, DARK, first=True)
# results table
rows, cols = 4, 3
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.7), Inches(2.7), Inches(6.0), Inches(2.6))
tbl = tbl_shape.table
hdr = ["Model", "F1", "ROC-AUC"]
data = [["Logistic Regression", "0.544", "0.879"],
        ["Gradient Boosting", "0.675", "0.927"],
        ["Random Forest  ★", "0.783", "0.966"]]
for j, htext in enumerate(hdr):
    cell = tbl.cell(0, j); cell.text = htext
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.color.rgb = WHITE
    p.runs[0].font.bold = True; p.runs[0].font.size = Pt(16)
for i, row in enumerate(data, start=1):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j); cell.text = val
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(15)
        win = (i == 3)
        p.runs[0].font.bold = win
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE3, 0xF3, 0xF2) if win else WHITE
        p.runs[0].font.color.rgb = TEAL if win else DARK
tf2 = box(s, 0.7, 5.5, 6.0, 1.2)
para(tf2, "Random Forest wins — ROC-AUC 0.966. AUC is the key metric because the classes are "
          "imbalanced.", 16, DARK, first=True)
# ROC image
s.shapes.add_picture("reports/figures/roc_curve.png", Inches(7.0), Inches(1.6),
                     height=Inches(5.3))

# ---------- Slide 5: Tuning + Registry ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE); title_bar(s, "Tuning & the Model Registry")
tf = box(s, 0.7, 1.6, 5.9, 5.3)
para(tf, "Hyperparameter tuning — Optuna", 22, TEAL, bold=True, first=True)
para(tf, "20 automated trials searching n_estimators, learning rate, depth", 17, DARK, bullet=True)
para(tf, "Every trial logged as a nested MLflow run — full audit trail", 17, DARK, bullet=True)
para(tf, "Best tuned model: ROC-AUC ≈ 0.958", 17, DARK, bullet=True, space=18)
para(tf, "Model Registry — lifecycle control", 22, TEAL, bold=True)
para(tf, "Best run registered as ecommerce-churn-model", 17, DARK, bullet=True)
para(tf, "Promoted: new version → Staging", 17, DARK, bullet=True)
para(tf, "Quality gate: ROC-AUC ≥ 0.85 → promoted to Production", 17, DARK, bullet=True)
para(tf, "The gate stops a bad model from ever going live", 16, GREY, space=4)
# confusion matrix
s.shapes.add_picture("reports/figures/confusion_matrix.png", Inches(7.0), Inches(1.8),
                     height=Inches(4.9))

# ---------- Slide 6: Serving + Monitoring + Airflow ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE); title_bar(s, "Deployment, Monitoring & Orchestration")
tf = box(s, 0.7, 1.6, 12, 5.3)
para(tf, "Deployment — real-time serving", 22, TEAL, bold=True, first=True)
para(tf, "Production model served over HTTP; send a customer record → get a churn prediction back",
     17, DARK, bullet=True)
para(tf, "Demo: a short-tenure customer who filed a complaint → predicted to churn (1)", 16, GREY, space=16)
para(tf, "Monitoring — drift detection", 22, TEAL, bold=True)
para(tf, "Simulated incoming batches; live metrics logged to MLflow", 17, DARK, bullet=True)
para(tf, "Population Stability Index (PSI) per feature; PSI > 0.2 = drift → would trigger a retrain",
     17, DARK, bullet=True, space=16)
para(tf, "Orchestration — Apache Airflow", 22, TEAL, bold=True)
para(tf, "A DAG runs ingest → train → tune → register on a schedule, calling the same code", 17, DARK, bullet=True)

# ---------- Slide 7: Conclusion ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
tf = box(s, 1, 1.3, 11.3, 1.2)
para(tf, "Key Takeaways", 36, WHITE, bold=True, first=True)
tf2 = box(s, 1, 2.7, 11.3, 4)
para(tf2, "MLflow let me manage the complete ML lifecycle in one reproducible system.", 22, TEAL, bold=True, first=True, space=16)
para(tf2, "Tracking made every experiment comparable and reproducible.", 19, WHITE, bullet=True)
para(tf2, "The registry gave a controlled Staging → Production path with a quality gate.", 19, WHITE, bullet=True)
para(tf2, "The same logged model served both batch and real-time predictions.", 19, WHITE, bullet=True)
para(tf2, "Monitoring closes the loop — drift detection signals when to retrain.", 19, WHITE, bullet=True)
para(tf2, "Best model: Random Forest, ROC-AUC 0.966, live in Production.", 19, TEAL, bullet=True, space=20)
para(tf2, "Thank you!  —  Demo: http://127.0.0.1:8080", 20, WHITE, bold=True)

prs.save("reports/Churn_MLOps_Presentation.pptx")
print("Saved reports/Churn_MLOps_Presentation.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
